"""Document text extraction utilities for nanobot."""

import mimetypes
from pathlib import Path

from loguru import logger


# Supported file extensions for text extraction
SUPPORTED_EXTENSIONS: set[str] = {
    ".pdf", ".docx", ".xlsx", ".pptx",
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
    ".log", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    ".png", ".jpg", ".jpeg", ".gif", ".webp",
}

_MAX_TEXT_LENGTH = 200_000


def extract_text(path: Path) -> str | None:
    """Extract text from a file."""
    if not isinstance(path, Path):
        path = Path(path)
    if not path.exists():
        return f"[error: file not found: {path}]"

    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf(path)
    elif ext == ".docx":
        return _extract_docx(path)
    elif ext == ".xlsx":
        return _extract_xlsx(path)
    elif ext == ".pptx":
        return _extract_pptx(path)
    elif _is_text_extension(ext):
        return _extract_text_file(path)
    elif ext in {".png", ".jpg", ".jpeg", ".gif", ".webp"}:
        return f"[image: {path.name}]"
    else:
        return None


def _extract_pdf(path: Path) -> str:
    try:
        from pypdf import PdfReader
    except ImportError:
        return "[error: pypdf not installed]"
    try:
        reader = PdfReader(path)
        pages = []
        for i, page in enumerate(reader.pages, 1):
            text = page.extract_text() or ""
            pages.append(f"--- Page {i} ---\n{text}")
        return _truncate("\n\n".join(pages), _MAX_TEXT_LENGTH)
    except Exception as e:
        return f"[error: failed to extract PDF: {e!s}]"


def _extract_docx(path: Path) -> str:
    try:
        from docx import Document as DocxDocument
    except ImportError:
        return "[error: python-docx not installed]"
    try:
        doc = DocxDocument(path)
        paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
        return _truncate("\n\n".join(paragraphs), _MAX_TEXT_LENGTH)
    except Exception as e:
        return f"[error: failed to extract DOCX: {e!s}]"


def _extract_xlsx(path: Path) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[error: openpyxl not installed]"
    try:
        wb = load_workbook(path, read_only=True, data_only=True)
        try:
            sheets = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                rows = []
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        rows.append(row_text)
                if rows:
                    sheets.append(f"--- Sheet: {sheet_name} ---\n" + "\n".join(rows))
            return _truncate("\n\n".join(sheets), _MAX_TEXT_LENGTH)
        finally:
            wb.close()
    except Exception as e:
        return f"[error: failed to extract XLSX: {e!s}]"


def _extract_pptx(path: Path) -> str:
    try:
        from pptx import Presentation as PptxPresentation
    except ImportError:
        return "[error: python-pptx not installed]"
    try:
        prs = PptxPresentation(path)
        slides = []
        for i, slide in enumerate(prs.slides, 1):
            slide_text = []
            for shape in slide.shapes:
                _collect_pptx_shape_text(shape, slide_text)
            if slide_text:
                slides.append(f"--- Slide {i} ---\n" + "\n".join(slide_text))
        return _truncate("\n\n".join(slides), _MAX_TEXT_LENGTH)
    except Exception as e:
        return f"[error: failed to extract PPTX: {e!s}]"


def _collect_pptx_shape_text(shape, out: list[str]) -> None:
    sub_shapes = getattr(shape, "shapes", None)
    if sub_shapes is not None:
        for sub in sub_shapes:
            _collect_pptx_shape_text(sub, out)
        return
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            cells = [cell.text.strip() for cell in row.cells]
            line = "\t".join(cell for cell in cells if cell)
            if line:
                out.append(line)
        return
    text = getattr(shape, "text", "")
    if text:
        out.append(text)


def _extract_text_file(path: Path) -> str:
    try:
        try:
            content = path.read_text(encoding="utf-8")
        except UnicodeDecodeError:
            content = path.read_text(encoding="latin-1")
        return _truncate(content, _MAX_TEXT_LENGTH)
    except Exception as e:
        return f"[error: failed to read file: {e!s}]"


def _truncate(text: str, max_length: int) -> str:
    if len(text) <= max_length:
        return text
    return text[:max_length] + f"... (truncated, {len(text)} chars total)"


def _is_text_extension(ext: str) -> bool:
    return ext in {
        ".txt", ".md", ".csv", ".json", ".xml", ".html", ".htm",
        ".log", ".yaml", ".yml", ".toml", ".ini", ".cfg",
    }


# ── Image detection ──

def detect_image_mime(data: bytes) -> str | None:
    """Detect image MIME type from the first bytes of a file."""
    if data[:8] == b"\x89PNG\r\n\x1a\n":
        return "image/png"
    if data[:2] in (b"\xff\xd8",):
        return "image/jpeg"
    if data[:6] in (b"GIF87a", b"GIF89a"):
        return "image/gif"
    if data[:4] == b"RIFF" and data[8:12] == b"WEBP":
        return "image/webp"
    return None


def is_image_file(path: str) -> bool:
    """Check whether a path looks like an image file."""
    p = Path(path)
    mime = None
    if p.is_file():
        try:
            with p.open("rb") as f:
                mime = detect_image_mime(f.read(16))
        except OSError:
            pass
    if not mime:
        mime = mimetypes.guess_type(path)[0]
    return bool(mime and mime.startswith("image/"))


def image_placeholder_text(path: str | None = None, *, empty: str = "[image]") -> str:
    """Generate placeholder text for an image."""
    if path:
        name = Path(path).name
        return f"[image: {name}]"
    return empty


def reference_non_image_attachments(content: str, media: list[str]) -> tuple[str, list[str]]:
    """Separate images from non-image attachments, just referencing them."""
    image_paths = []
    attachment_refs = []
    for path in media:
        if is_image_file(path):
            image_paths.append(path)
        else:
            attachment_refs.append(f"[Attachment: {path}]")
    if attachment_refs:
        suffix = "\n".join(attachment_refs)
        content = f"{content}\n\n{suffix}" if content else suffix
    return content, image_paths


_MAX_EXTRACT_FILE_SIZE = 50 * 1024 * 1024  # 50 MB


def extract_documents(
    text: str,
    media_paths: list[str],
    *,
    max_file_size: int = _MAX_EXTRACT_FILE_SIZE,
) -> tuple[str, list[str]]:
    """Separate images from documents in media_paths.

    Documents have their text extracted and appended to text.
    Only image paths are kept in the returned list.
    """
    image_paths = []
    doc_texts = []

    for path_str in media_paths:
        p = Path(path_str)
        if not p.is_file():
            continue

        try:
            size = p.stat().st_size
        except OSError:
            continue
        if size > max_file_size:
            logger.warning(
                "Skipping oversized file: {} ({:.1f} MB > {} MB limit)",
                p.name, size / (1024 * 1024), max_file_size // (1024 * 1024),
            )
            continue

        if is_image_file(path_str):
            image_paths.append(path_str)
        else:
            extracted = extract_text(p)
            if extracted and not extracted.startswith("[error:"):
                doc_texts.append(f"[File: {p.name}]\n{extracted}")

    if doc_texts:
        text = text + "\n\n" + "\n\n".join(doc_texts)

    return text, image_paths
